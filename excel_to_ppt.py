import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
import os

# Input/Output files (defaults)
excel_file = "data.xlsx"
template_ppt = "PPT Master Template.pptx"
output_ppt = "excel_barcharts.pptx"

BAR_COLOUR = RGBColor(0x00, 0xBC, 0xF2)

# Color palette for multi-series
SERIES_COLORS = [
    RGBColor(0x00, 0xBC, 0xF2),  # Blue
    RGBColor(0xFF, 0x6B, 0x6B),  # Red
    RGBColor(0x4E, 0xCB, 0x71),  # Green
    RGBColor(0xFF, 0xB8, 0x4D),  # Orange
    RGBColor(0x95, 0x5F, 0xE6),  # Purple
    RGBColor(0xF7, 0xCA, 0x18),  # Yellow
    RGBColor(0x3D, 0xC1, 0xD3),  # Cyan
    RGBColor(0xE7, 0x4C, 0x3C),  # Dark Red
]

# Available chart types
CHART_TYPES = {
    "Bar Chart": XL_CHART_TYPE.BAR_CLUSTERED,
    "Column Chart": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "Pie Chart": XL_CHART_TYPE.PIE,
    "Line Chart": XL_CHART_TYPE.LINE,
    "Area Chart": XL_CHART_TYPE.AREA,
    "Doughnut Chart": XL_CHART_TYPE.DOUGHNUT,
    "Stacked Bar": XL_CHART_TYPE.BAR_STACKED,
    "Stacked Column": XL_CHART_TYPE.COLUMN_STACKED,
}

class ChartConfigUI:
    def __init__(self, root):
        self.root = root
        self.root.title("PowerPoint Chart Generator - Multi-Series Edition")
        self.root.geometry("1100x750")
        self.root.resizable(True, True)
        
        # Variables
        self.excel_path = tk.StringVar(value=excel_file)
        self.template_path = tk.StringVar(value=template_ppt)
        self.output_path = tk.StringVar(value=output_ppt)
        self.starting_slide = tk.IntVar(value=3)
        
        # Chart selections and sheet info
        self.chart_selections = {}
        self.sheet_enabled = {}
        self.column_selections = {}  # Now stores list of selected columns
        self.percentage_mode = {}
        self.valid_sheets = []
        self.all_sheets_info = []
        
        self.setup_ui()
        self.load_excel_info()
        
    def setup_ui(self):
        # Main frame with scrollbar
        canvas = tk.Canvas(self.root)
        scrollbar = ttk.Scrollbar(self.root, orient="vertical", command=canvas.yview)
        self.scrollable_frame = ttk.Frame(canvas)
        
        self.scrollable_frame.bind(
            "<Configure>",
            lambda e: canvas.configure(scrollregion=canvas.bbox("all"))
        )
        
        canvas.create_window((0, 0), window=self.scrollable_frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)
        
        # Pack scrollable elements
        canvas.pack(side="left", fill="both", expand=True)
        scrollbar.pack(side="right", fill="y")
        
        # Main content frame
        main_frame = ttk.Frame(self.scrollable_frame, padding="10")
        main_frame.pack(fill=tk.BOTH, expand=True)
        main_frame.columnconfigure(1, weight=1)
        
        row = 0
        
        # Title
        title_label = ttk.Label(main_frame, text="PowerPoint Chart Generator - Multi-Series Edition", 
                               font=('Arial', 16, 'bold'))
        title_label.grid(row=row, column=0, columnspan=3, pady=(0, 20))
        row += 1
        
        # File selection section
        files_frame = ttk.LabelFrame(main_frame, text="File Configuration", padding="10")
        files_frame.grid(row=row, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=(0, 10))
        files_frame.columnconfigure(1, weight=1)
        row += 1
        
        # Excel file
        ttk.Label(files_frame, text="Excel File:").grid(row=0, column=0, sticky=tk.W, padx=(0, 5))
        ttk.Entry(files_frame, textvariable=self.excel_path, width=60).grid(row=0, column=1, sticky=(tk.W, tk.E), padx=(0, 5))
        ttk.Button(files_frame, text="Browse", command=self.browse_excel).grid(row=0, column=2)
        
        # Template file
        ttk.Label(files_frame, text="Template PPT:").grid(row=1, column=0, sticky=tk.W, padx=(0, 5), pady=(5, 0))
        ttk.Entry(files_frame, textvariable=self.template_path, width=60).grid(row=1, column=1, sticky=(tk.W, tk.E), padx=(0, 5), pady=(5, 0))
        ttk.Button(files_frame, text="Browse", command=self.browse_template).grid(row=1, column=2, pady=(5, 0))
        
        # Output file
        ttk.Label(files_frame, text="Output PPT:").grid(row=2, column=0, sticky=tk.W, padx=(0, 5), pady=(5, 0))
        ttk.Entry(files_frame, textvariable=self.output_path, width=60).grid(row=2, column=1, sticky=(tk.W, tk.E), padx=(0, 5), pady=(5, 0))
        ttk.Button(files_frame, text="Save As", command=self.browse_output).grid(row=2, column=2, pady=(5, 0))
        
        # Starting slide configuration
        config_frame = ttk.Frame(files_frame)
        config_frame.grid(row=3, column=0, columnspan=3, pady=(10, 0), sticky=tk.W)
        
        ttk.Label(config_frame, text="Starting Slide Number:").pack(side=tk.LEFT, padx=(0, 5))
        ttk.Spinbox(config_frame, from_=1, to=100, textvariable=self.starting_slide, width=5).pack(side=tk.LEFT, padx=(0, 10))
        ttk.Button(config_frame, text="🔄 Refresh Excel Data", command=self.load_excel_info).pack(side=tk.LEFT, padx=10)
        
        # Info section
        self.info_frame = ttk.LabelFrame(main_frame, text="Excel Data Analysis", padding="10")
        self.info_frame.grid(row=row, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=(0, 10))
        row += 1
        
        self.info_label = ttk.Label(self.info_frame, text="Load Excel file to analyze available data...")
        self.info_label.grid(row=0, column=0, sticky=tk.W)
        
        # Chart selection section
        self.chart_frame = ttk.LabelFrame(main_frame, text="Chart Configuration for Each Sheet", padding="10")
        self.chart_frame.grid(row=row, column=0, columnspan=3, sticky=(tk.W, tk.E, tk.N, tk.S), pady=(0, 10))
        self.chart_frame.columnconfigure(0, weight=1)
        main_frame.rowconfigure(row, weight=1)
        row += 1
        
        # Quick selection section
        quick_frame = ttk.LabelFrame(main_frame, text="Batch Operations", padding="10")
        quick_frame.grid(row=row, column=0, columnspan=3, sticky=(tk.W, tk.E), pady=(0, 10))
        row += 1
        
        # Quick selection buttons row 1
        quick_row1 = ttk.Frame(quick_frame)
        quick_row1.pack(fill=tk.X, pady=2)
        
        ttk.Button(quick_row1, text="✓ Enable All Sheets", command=self.enable_all_sheets).pack(side=tk.LEFT, padx=2)
        ttk.Button(quick_row1, text="✗ Disable All Sheets", command=self.disable_all_sheets).pack(side=tk.LEFT, padx=2)
        ttk.Button(quick_row1, text="📊 All Bar Charts", command=lambda: self.set_all_charts("Bar Chart")).pack(side=tk.LEFT, padx=2)
        ttk.Button(quick_row1, text="📈 All Column Charts", command=lambda: self.set_all_charts("Column Chart")).pack(side=tk.LEFT, padx=2)
        
        # Quick selection buttons row 2
        quick_row2 = ttk.Frame(quick_frame)
        quick_row2.pack(fill=tk.X, pady=2)
        
        ttk.Button(quick_row2, text="🥧 All Pie Charts", command=lambda: self.set_all_charts("Pie Chart")).pack(side=tk.LEFT, padx=2)
        ttk.Button(quick_row2, text="📉 All Line Charts", command=lambda: self.set_all_charts("Line Chart")).pack(side=tk.LEFT, padx=2)
        ttk.Button(quick_row2, text="% Enable All Percentage Mode", command=self.enable_all_percentage).pack(side=tk.LEFT, padx=2)
        ttk.Button(quick_row2, text="# Disable All Percentage Mode", command=self.disable_all_percentage).pack(side=tk.LEFT, padx=2)
        
        # Action buttons
        button_frame = ttk.Frame(main_frame)
        button_frame.grid(row=row, column=0, columnspan=3, pady=(20, 0))
        
        ttk.Button(button_frame, text="🚀 Generate PowerPoint", command=self.generate_ppt, 
                  style="Accent.TButton").pack(side=tk.LEFT, padx=5)
        ttk.Button(button_frame, text="❌ Exit", command=self.root.quit).pack(side=tk.LEFT, padx=5)
        
        # Bind mousewheel to canvas
        def _on_mousewheel(event):
            canvas.yview_scroll(int(-1*(event.delta/120)), "units")
        canvas.bind_all("<MouseWheel>", _on_mousewheel)
        
    def browse_excel(self):
        filename = filedialog.askopenfilename(
            title="Select Excel File",
            filetypes=[("Excel files", "*.xlsx *.xls"), ("All files", "*.*")]
        )
        if filename:
            self.excel_path.set(filename)
            self.load_excel_info()
    
    def browse_template(self):
        filename = filedialog.askopenfilename(
            title="Select PowerPoint Template",
            filetypes=[("PowerPoint files", "*.pptx *.ppt"), ("All files", "*.*")]
        )
        if filename:
            self.template_path.set(filename)
    
    def browse_output(self):
        filename = filedialog.asksaveasfilename(
            title="Save PowerPoint As",
            filetypes=[("PowerPoint files", "*.pptx")],
            defaultextension=".pptx"
        )
        if filename:
            self.output_path.set(filename)
    
    def load_excel_info(self):
        try:
            if not os.path.exists(self.excel_path.get()):
                self.info_label.config(text="❌ Excel file not found")
                return
            
            # Load all Excel sheets
            sheets = pd.read_excel(self.excel_path.get(), sheet_name=None, header=2)
            
            # Analyze all sheets
            self.all_sheets_info = []
            self.valid_sheets = []
            
            for sheet_name, df in sheets.items():
                sheet_info = {
                    'name': sheet_name,
                    'total_rows': len(df),
                    'total_columns': len(df.columns),
                    'is_valid': False,
                    'valid_rows': 0,
                    'has_numeric_data': False,
                    'column_names': list(df.columns) if not df.empty else [],
                    'numeric_columns': []
                }
                
                if not df.empty and len(df.columns) >= 2:
                    # Test data cleaning process
                    test_df = df.copy()
                    test_df = test_df[~test_df.iloc[:, 0].astype(str).str.startswith("Base")]
                    test_df = test_df.dropna(subset=[test_df.columns[0]])
                    
                    # Find all numeric columns (skip first column which is labels)
                    numeric_columns = []
                    for col_idx in range(1, len(df.columns)):
                        col_data = pd.to_numeric(test_df.iloc[:, col_idx], errors='coerce')
                        if col_data.notna().sum() > 0:
                            numeric_columns.append({
                                'index': col_idx,
                                'name': str(df.columns[col_idx]),
                                'valid_count': col_data.notna().sum()
                            })
                    
                    if numeric_columns and not test_df.empty:
                        sheet_info['is_valid'] = True
                        sheet_info['valid_rows'] = len(test_df)
                        sheet_info['has_numeric_data'] = True
                        sheet_info['numeric_columns'] = numeric_columns
                        self.valid_sheets.append(sheet_name)
                
                self.all_sheets_info.append(sheet_info)
            
            # Update info display
            total_sheets = len(sheets)
            valid_sheets_count = len(self.valid_sheets)
            
            info_text = f"📊 Excel Analysis Results:\n"
            info_text += f"   • Total sheets found: {total_sheets}\n"
            info_text += f"   • Sheets with valid chart data: {valid_sheets_count}\n"
            info_text += f"   • Charts will start from slide: {self.starting_slide.get()}\n"
            
            if valid_sheets_count > 0:
                info_text += f"   • Valid sheets: {', '.join(self.valid_sheets[:5])}"
                if len(self.valid_sheets) > 5:
                    info_text += f" and {len(self.valid_sheets)-5} more..."
            
            self.info_label.config(text=info_text)
            
            # Create dynamic chart selectors
            self.create_dynamic_selectors()
            
        except Exception as e:
            self.info_label.config(text=f"❌ Error analyzing Excel file: {str(e)}")
    
    def open_series_selector(self, sheet_name, sheet_info):
        """Open a dialog to select multiple series for a sheet"""
        dialog = tk.Toplevel(self.root)
        dialog.title(f"Select Data Series - {sheet_name}")
        dialog.geometry("500x400")
        dialog.transient(self.root)
        dialog.grab_set()
        
        # Info label
        info_label = ttk.Label(dialog, text=f"Select one or more data columns to chart from '{sheet_name}':", 
                              font=('Arial', 10, 'bold'))
        info_label.pack(pady=10, padx=10)
        
        # Listbox with scrollbar for column selection
        list_frame = ttk.Frame(dialog)
        list_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=5)
        
        listbox_scroll = ttk.Scrollbar(list_frame, orient="vertical")
        listbox = tk.Listbox(list_frame, selectmode=tk.MULTIPLE, 
                            yscrollcommand=listbox_scroll.set, height=15)
        listbox_scroll.config(command=listbox.yview)
        
        listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        listbox_scroll.pack(side=tk.RIGHT, fill=tk.Y)
        
        # Populate listbox with numeric columns
        for col in sheet_info['numeric_columns']:
            listbox.insert(tk.END, col['name'])
        
        # Pre-select previously selected columns if they exist
        if sheet_name in self.column_selections and self.column_selections[sheet_name]:
            for i, col in enumerate(sheet_info['numeric_columns']):
                if col['index'] in self.column_selections[sheet_name]['indices']:
                    listbox.selection_set(i)
        else:
            # Default: select first column
            listbox.selection_set(0)
        
        # Helper buttons
        helper_frame = ttk.Frame(dialog)
        helper_frame.pack(pady=5)
        
        ttk.Button(helper_frame, text="Select All", 
                  command=lambda: listbox.selection_set(0, tk.END)).pack(side=tk.LEFT, padx=5)
        ttk.Button(helper_frame, text="Clear All", 
                  command=lambda: listbox.selection_clear(0, tk.END)).pack(side=tk.LEFT, padx=5)
        
        # Selection info
        selection_label = ttk.Label(dialog, text="", foreground="gray")
        selection_label.pack(pady=5)
        
        def update_selection_info():
            count = len(listbox.curselection())
            if count == 0:
                selection_label.config(text="⚠️ Please select at least one column")
            elif count == 1:
                selection_label.config(text="✓ Single series chart")
            else:
                selection_label.config(text=f"✓ Multi-series chart ({count} series)")
        
        listbox.bind('<<ListboxSelect>>', lambda e: update_selection_info())
        update_selection_info()
        
        # Buttons
        button_frame = ttk.Frame(dialog)
        button_frame.pack(pady=10)
        
        def save_selection():
            selected_indices = listbox.curselection()
            if not selected_indices:
                messagebox.showwarning("No Selection", "Please select at least one column!")
                return
            
            # Save the selected column indices and names
            selected_columns = []
            selected_names = []
            for idx in selected_indices:
                col = sheet_info['numeric_columns'][idx]
                selected_columns.append(col['index'])
                selected_names.append(col['name'])
            
            self.column_selections[sheet_name] = {
                'indices': selected_columns,
                'names': selected_names,
                'columns': sheet_info['numeric_columns']
            }
            
            # Update the button text
            self.update_series_button_text(sheet_name)
            
            dialog.destroy()
        
        ttk.Button(button_frame, text="✓ Save Selection", 
                  command=save_selection).pack(side=tk.LEFT, padx=5)
        ttk.Button(button_frame, text="Cancel", 
                  command=dialog.destroy).pack(side=tk.LEFT, padx=5)
    
    def update_series_button_text(self, sheet_name):
        """Update the series selection button text to show what's selected"""
        if sheet_name in self.column_selections and self.column_selections[sheet_name]:
            count = len(self.column_selections[sheet_name]['indices'])
            names = self.column_selections[sheet_name]['names']
            
            if count == 1:
                text = names[0][:15]
            else:
                text = f"{count} series"
            
            # Find the button and update its text
            for widget in self.chart_frame.winfo_children():
                if isinstance(widget, tk.Canvas):
                    scrollable_frame = widget.winfo_children()[0] if widget.winfo_children() else None
                    if scrollable_frame:
                        for frame in scrollable_frame.winfo_children():
                            if hasattr(frame, 'sheet_name') and frame.sheet_name == sheet_name:
                                if hasattr(frame, 'series_button'):
                                    frame.series_button.config(text=text)
    
    def create_dynamic_selectors(self):
        # Clear existing selectors
        for widget in self.chart_frame.winfo_children():
            widget.destroy()
        
        self.chart_selections.clear()
        self.sheet_enabled.clear()
        # Don't clear column_selections - preserve previous selections
        self.percentage_mode.clear()
        
        if not self.all_sheets_info:
            ttk.Label(self.chart_frame, text="No sheets found in Excel file").pack(pady=20)
            return
        
        # Create scrollable frame for selectors
        selector_canvas = tk.Canvas(self.chart_frame, height=400)
        selector_scrollbar = ttk.Scrollbar(self.chart_frame, orient="vertical", command=selector_canvas.yview)
        selector_scrollable_frame = ttk.Frame(selector_canvas)
        
        selector_scrollable_frame.bind(
            "<Configure>",
            lambda e: selector_canvas.configure(scrollregion=selector_canvas.bbox("all"))
        )
        
        selector_canvas.create_window((0, 0), window=selector_scrollable_frame, anchor="nw")
        selector_canvas.configure(yscrollcommand=selector_scrollbar.set)
        
        # Header
        header_frame = ttk.Frame(selector_scrollable_frame)
        header_frame.pack(fill=tk.X, padx=5, pady=5)
        
        ttk.Label(header_frame, text="✓", font=('Arial', 9, 'bold'), width=3).grid(row=0, column=0)
        ttk.Label(header_frame, text="Sheet Name", font=('Arial', 9, 'bold'), width=18).grid(row=0, column=1, sticky=tk.W, padx=5)
        ttk.Label(header_frame, text="Data Series", font=('Arial', 9, 'bold'), width=15).grid(row=0, column=2)
        ttk.Label(header_frame, text="%", font=('Arial', 9, 'bold'), width=3).grid(row=0, column=3)
        ttk.Label(header_frame, text="Chart Type", font=('Arial', 9, 'bold'), width=15).grid(row=0, column=4)
        ttk.Label(header_frame, text="Slide", font=('Arial', 9, 'bold'), width=6).grid(row=0, column=5)
        
        # Create selector for each sheet
        for i, sheet_info in enumerate(self.all_sheets_info):
            sheet_name = sheet_info['name']
            
            frame = ttk.Frame(selector_scrollable_frame)
            frame.pack(fill=tk.X, padx=5, pady=2)
            
            # Enable/Disable checkbox
            enabled_var = tk.BooleanVar(value=sheet_info['is_valid'])
            self.sheet_enabled[sheet_name] = enabled_var
            
            checkbox = ttk.Checkbutton(frame, variable=enabled_var, command=self.update_slide_numbers)
            checkbox.grid(row=0, column=0, padx=5)
            
            # Sheet name with status indicator
            name_text = sheet_name[:16] + "..." if len(sheet_name) > 16 else sheet_name
            if sheet_info['is_valid']:
                name_text = f"✅ {name_text}"
            else:
                name_text = f"❌ {name_text}"
                enabled_var.set(False)
            
            name_label = ttk.Label(frame, text=name_text, width=20)
            name_label.grid(row=0, column=1, sticky=tk.W, padx=5)
            
            # Series selector button (NEW - replaces dropdown)
            if sheet_info['is_valid'] and sheet_info['numeric_columns']:
                # Initialize with first column if not already set
                if sheet_name not in self.column_selections or not self.column_selections[sheet_name]:
                    first_col = sheet_info['numeric_columns'][0]
                    self.column_selections[sheet_name] = {
                        'indices': [first_col['index']],
                        'names': [first_col['name']],
                        'columns': sheet_info['numeric_columns']
                    }
                
                # Determine button text
                count = len(self.column_selections[sheet_name]['indices'])
                if count == 1:
                    button_text = self.column_selections[sheet_name]['names'][0][:13]
                else:
                    button_text = f"{count} series"
                
                series_button = ttk.Button(frame, text=button_text, width=15,
                                          command=lambda sn=sheet_name, si=sheet_info: self.open_series_selector(sn, si))
                series_button.grid(row=0, column=2, padx=5)
                frame.series_button = series_button
            else:
                self.column_selections[sheet_name] = None
                ttk.Label(frame, text="N/A", width=13).grid(row=0, column=2)
            
            # Percentage mode checkbox
            percentage_var = tk.BooleanVar(value=False)
            self.percentage_mode[sheet_name] = percentage_var
            
            if sheet_info['is_valid']:
                percentage_check = ttk.Checkbutton(frame, variable=percentage_var)
                percentage_check.grid(row=0, column=3, padx=5)
            else:
                ttk.Label(frame, text="", width=3).grid(row=0, column=3)
            
            # Chart type selector
            chart_var = tk.StringVar(value="Bar Chart")
            self.chart_selections[sheet_name] = chart_var
            
            if sheet_info['is_valid']:
                combo = ttk.Combobox(frame, textvariable=chart_var, values=list(CHART_TYPES.keys()), 
                                   state="readonly", width=13)
            else:
                combo = ttk.Combobox(frame, textvariable=chart_var, values=list(CHART_TYPES.keys()), 
                                   state="disabled", width=13)
            
            combo.grid(row=0, column=4, padx=5)
            
            # Slide number
            slide_label = ttk.Label(frame, text="", width=6)
            slide_label.grid(row=0, column=5)
            
            # Store references
            frame.slide_label = slide_label
            frame.sheet_name = sheet_name
        
        # Pack the canvas and scrollbar
        selector_canvas.pack(side="left", fill="both", expand=True)
        selector_scrollbar.pack(side="right", fill="y")
        
        # Update slide numbers
        self.update_slide_numbers()
        
        # Bind mousewheel to selector canvas
        def _on_selector_mousewheel(event):
            selector_canvas.yview_scroll(int(-1*(event.delta/120)), "units")
        selector_canvas.bind_all("<Button-4>", lambda e: selector_canvas.yview_scroll(-1, "units"))
        selector_canvas.bind_all("<Button-5>", lambda e: selector_canvas.yview_scroll(1, "units"))
    
    def update_slide_numbers(self):
        """Update slide numbers based on enabled sheets"""
        slide_num = self.starting_slide.get()
        
        for widget in self.chart_frame.winfo_children():
            if isinstance(widget, tk.Canvas):
                scrollable_frame = widget.winfo_children()[0] if widget.winfo_children() else None
                if scrollable_frame:
                    for frame in scrollable_frame.winfo_children():
                        if hasattr(frame, 'slide_label') and hasattr(frame, 'sheet_name'):
                            if self.sheet_enabled.get(frame.sheet_name, tk.BooleanVar()).get():
                                frame.slide_label.config(text=f"{slide_num}")
                                slide_num += 1
                            else:
                                frame.slide_label.config(text="—")
    
    def enable_all_sheets(self):
        for sheet_name, var in self.sheet_enabled.items():
            sheet_info = next((s for s in self.all_sheets_info if s['name'] == sheet_name), None)
            if sheet_info and sheet_info['is_valid']:
                var.set(True)
        self.update_slide_numbers()
    
    def disable_all_sheets(self):
        for var in self.sheet_enabled.values():
            var.set(False)
        self.update_slide_numbers()
    
    def enable_all_percentage(self):
        for sheet_name, var in self.percentage_mode.items():
            if self.sheet_enabled.get(sheet_name, tk.BooleanVar()).get():
                var.set(True)
    
    def disable_all_percentage(self):
        for var in self.percentage_mode.values():
            var.set(False)
    
    def set_all_charts(self, chart_type):
        for sheet_name, var in self.chart_selections.items():
            if self.sheet_enabled.get(sheet_name, tk.BooleanVar()).get():
                var.set(chart_type)
    
    def get_enabled_sheets(self):
        """Get list of enabled sheets with their configuration"""
        enabled_sheets = []
        slide_num = self.starting_slide.get()
        
        for sheet_info in self.all_sheets_info:
            sheet_name = sheet_info['name']
            if self.sheet_enabled.get(sheet_name, tk.BooleanVar()).get() and sheet_info['is_valid']:
                column_info = self.column_selections.get(sheet_name)
                
                if column_info and column_info['indices']:
                    enabled_sheets.append({
                        'name': sheet_name,
                        'chart_type': self.chart_selections[sheet_name].get(),
                        'slide_number': slide_num,
                        'data_rows': sheet_info['valid_rows'],
                        'column_indices': column_info['indices'],
                        'column_names': column_info['names'],
                        'percentage_mode': self.percentage_mode.get(sheet_name, tk.BooleanVar()).get()
                    })
                    slide_num += 1
        
        return enabled_sheets
    
    def generate_ppt(self):
        enabled_sheets = self.get_enabled_sheets()
        
        if not enabled_sheets:
            messagebox.showerror("Error", "No sheets selected for chart generation!")
            return
        
        try:
            progress_window = tk.Toplevel(self.root)
            progress_window.title("Generating PowerPoint...")
            progress_window.geometry("500x200")
            progress_window.transient(self.root)
            progress_window.grab_set()
            
            progress_label = ttk.Label(progress_window, text="Initializing...", font=('Arial', 10))
            progress_label.pack(pady=20)
            
            progress_bar = ttk.Progressbar(progress_window, mode='determinate', length=400)
            progress_bar.pack(pady=10)
            progress_bar['maximum'] = len(enabled_sheets) + 2
            
            detail_label = ttk.Label(progress_window, text="", font=('Arial', 9), foreground='gray')
            detail_label.pack(pady=5)
            
            self.root.update()
            
            self.create_powerpoint(progress_label, progress_bar, detail_label, enabled_sheets)
            
            progress_window.destroy()
            
            success_msg = f"PowerPoint created successfully!\n\n"
            success_msg += f"📊 Charts created: {len(enabled_sheets)}\n"
            success_msg += f"💾 Saved as: {os.path.basename(self.output_path.get())}\n\n"
            success_msg += f"📂 Full path: {self.output_path.get()}"
            
            messagebox.showinfo("Success! 🎉", success_msg)
            
        except Exception as e:
            messagebox.showerror("Error", f"Failed to create PowerPoint:\n\n{str(e)}")
    
    def create_powerpoint(self, progress_label, progress_bar, detail_label, enabled_sheets):
        progress_label.config(text="Loading Excel data...")
        detail_label.config(text=f"Reading {os.path.basename(self.excel_path.get())}")
        self.root.update()
        
        sheets = pd.read_excel(self.excel_path.get(), sheet_name=None)
        progress_bar['value'] += 1
        self.root.update()
        
        progress_label.config(text="Loading PowerPoint template...")
        detail_label.config(text=f"Opening {os.path.basename(self.template_path.get())}")
        self.root.update()
        
        prs = Presentation(self.template_path.get())
        slide_layout = prs.slide_layouts[min(2, len(prs.slide_layouts)-1)]
        progress_bar['value'] += 1
        self.root.update()
        
        # Process each enabled sheet
        for i, sheet_config in enumerate(enabled_sheets):
            sheet_name = sheet_config['name']
            chart_type_name = sheet_config['chart_type']
            chart_type = CHART_TYPES[chart_type_name]
            column_indices = sheet_config['column_indices']
            column_names = sheet_config['column_names']
            percentage_mode = sheet_config['percentage_mode']
            
            progress_label.config(text=f"Creating chart {i+1} of {len(enabled_sheets)}")
            series_info = f"{len(column_indices)} series" if len(column_indices) > 1 else "single series"
            detail_label.config(text=f"Processing {sheet_name} → {chart_type_name} ({series_info})")
            self.root.update()
            
            # Get and clean data
            df = sheets[sheet_name].copy()
            
            # Filter out Base rows and handle NaN in first column
            df = df[~df.iloc[:, 0].astype(str).str.startswith("Base")]
            df = df.dropna(subset=[df.columns[0]])
            
            # Extract the label column
            labels = df.iloc[:, 0].astype(str)
            
            # Create chart data with multiple series
            chart_data = CategoryChartData()
            chart_data.categories = labels.tolist()
            
            # Add each selected column as a series
            for col_idx, col_name in zip(column_indices, column_names):
                data_column = pd.to_numeric(df.iloc[:, col_idx], errors="coerce")
                
                # Create clean series data
                clean_data = data_column.fillna(0)  # Fill NaN with 0 for charting
                
                # Apply percentage conversion if enabled
                if percentage_mode:
                    clean_data = clean_data.round(1)
                
                chart_data.add_series(col_name, clean_data.tolist())
            
            # Create slide
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title
            pct_suffix = " (%)" if percentage_mode else ""
            if len(column_names) == 1:
                title_text = f"{sheet_name} - {column_names[0]}{pct_suffix}"
            else:
                title_text = f"{sheet_name} - Multi-Series Chart{pct_suffix}"
            
            if slide.shapes.title:
                slide.shapes.title.text = title_text
            else:
                title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
                title_frame = title_box.text_frame
                p = title_frame.paragraphs[0]
                p.text = title_text
                p.font.size = Pt(24)
            
            # Create chart
            x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(5)
            chart_shape = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
            chart = chart_shape.chart
            
            # Format chart
            self.format_chart(chart, chart_type, percentage_mode, len(column_indices))
            
            progress_bar['value'] += 1
            self.root.update()
        
        # Save presentation
        progress_label.config(text="Saving PowerPoint presentation...")
        detail_label.config(text=f"Writing to {os.path.basename(self.output_path.get())}")
        self.root.update()
        
        prs.save(self.output_path.get())
    
    def format_chart(self, chart, chart_type, percentage_mode=False, series_count=1):
        """Apply formatting based on chart type"""
        try:
            if chart_type in [XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT]:
                # Pie/Doughnut chart formatting
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.RIGHT
                chart.plots[0].has_data_labels = True
                
                if percentage_mode:
                    # Show custom percentage labels with 1 decimal
                    chart.plots[0].data_labels.show_percentage = False
                    chart.plots[0].data_labels.show_value = True
                    chart.plots[0].data_labels.number_format = '0.0"%"'
                else:
                    chart.plots[0].data_labels.show_percentage = True
                    chart.plots[0].data_labels.show_value = False
                
                chart.plots[0].data_labels.show_category_name = False
            else:
                # Bar, Column, Line, Area charts
                # Show legend if multi-series
                if series_count > 1:
                    chart.has_legend = True
                    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                    chart.legend.font.size = Pt(9)
                else:
                    chart.has_legend = False
                
                # Only set axis properties for charts that have axes
                if hasattr(chart, 'value_axis') and hasattr(chart, 'category_axis'):
                    try:
                        chart.value_axis.has_major_gridlines = False
                        chart.value_axis.has_minor_gridlines = False
                        chart.category_axis.has_major_gridlines = False
                        chart.category_axis.has_minor_gridlines = False
                        chart.value_axis.tick_labels.font.size = Pt(10)
                        chart.category_axis.tick_labels.font.size = Pt(10)
                        
                        # Set axis number format for percentage mode
                        if percentage_mode:
                            chart.value_axis.tick_labels.number_format = '0.0"%"'
                    except:
                        pass
                
                # Add data labels (optional for multi-series to avoid clutter)
                try:
                    if series_count == 1:
                        chart.plots[0].has_data_labels = True
                        chart.plots[0].data_labels.show_value = True
                        
                        # Set data label format for percentage mode
                        if percentage_mode:
                            chart.plots[0].data_labels.number_format = '0.0"%"'
                        else:
                            chart.plots[0].data_labels.number_format = '0.0'
                except:
                    pass
            
            # Set colors for all series
            for series_idx, series in enumerate(chart.series):
                try:
                    color = SERIES_COLORS[series_idx % len(SERIES_COLORS)]
                    
                    if chart_type in [XL_CHART_TYPE.PIE, XL_CHART_TYPE.DOUGHNUT]:
                        # For pie charts, color individual points
                        for point_idx, point in enumerate(series.points):
                            point_color = SERIES_COLORS[point_idx % len(SERIES_COLORS)]
                            point.format.fill.solid()
                            point.format.fill.fore_color.rgb = point_color
                    else:
                        # For other chart types
                        series.format.fill.solid()
                        series.format.fill.fore_color.rgb = color
                        
                        # Set line color for line charts
                        if chart_type == XL_CHART_TYPE.LINE:
                            series.format.line.color.rgb = color
                            series.format.line.width = Pt(2.5)
                        
                        # Set gap width for bar/column charts
                        if hasattr(chart.plots[0], 'gap_width') and chart_type in [
                            XL_CHART_TYPE.BAR_CLUSTERED, XL_CHART_TYPE.COLUMN_CLUSTERED,
                            XL_CHART_TYPE.BAR_STACKED, XL_CHART_TYPE.COLUMN_STACKED
                        ]:
                            chart.plots[0].gap_width = 50
                except Exception as e:
                    print(f"Warning: Could not set color for series {series_idx}: {e}")
            
            # Set data label font size
            try:
                chart.plots[0].data_labels.font.size = Pt(9)
            except:
                pass
                
        except Exception as e:
            print(f"Warning: Could not apply all formatting to chart: {e}")

def main():
    root = tk.Tk()
    app = ChartConfigUI(root)
    
    # Center the window
    root.update_idletasks()
    x = (root.winfo_screenwidth() // 2) - (root.winfo_width() // 2)
    y = (root.winfo_screenheight() // 2) - (root.winfo_height() // 2)
    root.geometry(f"+{x}+{y}")
    
    root.mainloop()

if __name__ == "__main__":
    main()